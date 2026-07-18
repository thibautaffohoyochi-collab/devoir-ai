"""
Extraction du contenu des fichiers uploadés (PDF, Word, PowerPoint, TXT)
"""
import io
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from fastapi import UploadFile, HTTPException

ALLOWED_TYPES = {
    "application/pdf": "pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
    "application/msword": "doc",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
    "text/plain": "txt",
}

MAX_SIZE_BYTES = 10 * 1024 * 1024  # 10 MB

async def extract_text(file: UploadFile) -> tuple[str, str]:
    """
    Retourne (texte_extrait, type_fichier)
    """
    if file.content_type not in ALLOWED_TYPES:
        raise HTTPException(
            status_code=400,
            detail=f"Type de fichier non supporté: {file.content_type}. Acceptés: PDF, DOCX, PPTX, TXT"
        )

    content = await file.read()

    if len(content) > MAX_SIZE_BYTES:
        raise HTTPException(status_code=413, detail="Fichier trop volumineux (max 10 MB)")

    file_type = ALLOWED_TYPES[file.content_type]

    if file_type == "pdf":
        text = _extract_pdf(content)
    elif file_type in ("docx", "doc"):
        text = _extract_docx(content)
    elif file_type == "pptx":
        text = _extract_pptx(content)
    else:
        text = content.decode("utf-8", errors="replace")

    if not text.strip():
        raise HTTPException(status_code=422, detail="Impossible d'extraire du texte de ce fichier")

    return text.strip(), file_type


def _extract_pdf(content: bytes) -> str:
    reader = PdfReader(io.BytesIO(content))
    pages = []
    for page in reader.pages:
        extracted = page.extract_text()
        if extracted:
            pages.append(extracted)
    return "\n".join(pages)


def _extract_docx(content: bytes) -> str:
    doc = Document(io.BytesIO(content))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n".join(paragraphs)


def _extract_pptx(content: bytes) -> str:
    prs = Presentation(io.BytesIO(content))
    slides_text = []
    for i, slide in enumerate(prs.slides, 1):
        slide_content = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_content.append(shape.text)
        if slide_content:
            slides_text.append(f"--- Diapositive {i} ---\n" + "\n".join(slide_content))
    return "\n\n".join(slides_text)
