"""
Export amélioré : PDF pro, Word académique, PowerPoint avec thème visuel
"""
import io
import re
import os
from datetime import datetime
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import cm
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer, HRFlowable,
    PageBreak, Table, TableStyle
)
from reportlab.lib.enums import TA_JUSTIFY, TA_CENTER, TA_LEFT
from reportlab.lib import colors
from docx import Document
from docx.shared import Inches, Pt, RGBColor as DocxRGB, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

RESULTS_DIR = "results"
os.makedirs(RESULTS_DIR, exist_ok=True)

# Palette couleurs
PRIMARY = colors.HexColor('#2563eb')
SECONDARY = colors.HexColor('#1e40af')
ACCENT = colors.HexColor('#7c3aed')
DARK = colors.HexColor('#0f172a')
LIGHT_BG = colors.HexColor('#f8fafc')
TEXT_COLOR = colors.HexColor('#1e293b')
MUTED = colors.HexColor('#64748b')


def export_response(text: str, format_sortie: str, titre: str, user_id: int) -> str:
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    safe = re.sub(r'[^\w\s-]', '', titre)[:40].strip().replace(' ', '_')
    filename = f"{user_id}_{safe}_{timestamp}"

    if format_sortie == "pdf":
        path = os.path.join(RESULTS_DIR, f"{filename}.pdf")
        _export_pdf(text, path, titre)
    elif format_sortie == "docx":
        path = os.path.join(RESULTS_DIR, f"{filename}.docx")
        _export_docx(text, path, titre)
    elif format_sortie == "pptx":
        path = os.path.join(RESULTS_DIR, f"{filename}.pptx")
        _export_pptx(text, path, titre)
    else:
        path = os.path.join(RESULTS_DIR, f"{filename}.pdf")
        _export_pdf(text, path, titre)
    return path


# ── PDF PRO ───────────────────────────────────────────────────────────────────

def _make_pdf_styles():
    styles = getSampleStyleSheet()
    return {
        'title': ParagraphStyle(
            'Title', fontSize=18, fontName='Helvetica-Bold',
            textColor=PRIMARY, spaceAfter=8, alignment=TA_CENTER,
            leading=22, wordWrap='CJK'
        ),
        'subtitle': ParagraphStyle(
            'Subtitle', fontSize=10, fontName='Helvetica',
            textColor=MUTED, spaceAfter=16, alignment=TA_CENTER,
            leading=14
        ),
        'h1': ParagraphStyle(
            'H1', fontSize=14, fontName='Helvetica-Bold',
            textColor=SECONDARY, spaceBefore=20, spaceAfter=8,
            borderPadding=(0, 0, 4, 0)
        ),
        'h2': ParagraphStyle(
            'H2', fontSize=12, fontName='Helvetica-Bold',
            textColor=ACCENT, spaceBefore=14, spaceAfter=6
        ),
        'body': ParagraphStyle(
            'Body', fontSize=11, fontName='Helvetica',
            textColor=TEXT_COLOR, leading=18, spaceAfter=8,
            alignment=TA_JUSTIFY
        ),
        'footer': ParagraphStyle(
            'Footer', fontSize=8, fontName='Helvetica',
            textColor=MUTED, alignment=TA_CENTER
        ),
    }


def _add_page_number(canvas, doc):
    canvas.saveState()
    canvas.setFont('Helvetica', 8)
    canvas.setFillColor(MUTED)
    page_num = canvas.getPageNumber()
    canvas.drawCentredString(A4[0] / 2, 1.5 * cm, f"— {page_num} —")
    canvas.drawString(2.5 * cm, 1.5 * cm, "DevoirAI")
    canvas.restoreState()


def _export_pdf(text: str, path: str, titre: str):
    doc = SimpleDocTemplate(
        path, pagesize=A4,
        leftMargin=2.5*cm, rightMargin=2.5*cm,
        topMargin=3*cm, bottomMargin=2.5*cm
    )
    styles = _make_pdf_styles()
    story = []

    # En-tête — titre sur sa propre ligne, date en dessous
    story.append(Spacer(1, 0.5*cm))
    story.append(Paragraph(titre, styles['title']))
    story.append(Spacer(1, 0.3*cm))
    story.append(Paragraph(
        f"Généré le {datetime.now().strftime('%d/%m/%Y à %H:%M')}",
        styles['subtitle']
    ))
    story.append(Spacer(1, 0.3*cm))
    story.append(HRFlowable(width="100%", thickness=2, color=PRIMARY, spaceAfter=20))

    # Contenu
    for line in text.split('\n'):
        line = line.strip()
        if not line:
            story.append(Spacer(1, 0.3*cm))
            continue

        # Echapper les caractères XML
        safe = line.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')

        if line.startswith('# '):
            story.append(Paragraph(safe[2:], styles['h1']))
            story.append(HRFlowable(width="100%", thickness=0.5, color=PRIMARY, spaceAfter=4))
        elif line.startswith('## '):
            story.append(Paragraph(safe[3:], styles['h2']))
        elif line.startswith('### '):
            story.append(Paragraph(f"<b>{safe[4:]}</b>", styles['body']))
        elif line.startswith('**') and line.endswith('**'):
            story.append(Paragraph(f"<b>{safe[2:-2]}</b>", styles['body']))
        elif line.startswith('- ') or line.startswith('• '):
            story.append(Paragraph(f"&bull; {safe[2:]}", styles['body']))
        else:
            story.append(Paragraph(safe, styles['body']))

    # Pied de page
    story.append(Spacer(1, 1*cm))
    story.append(HRFlowable(width="100%", thickness=0.5, color=MUTED))
    story.append(Paragraph("Document généré par DevoirAI", styles['footer']))

    doc.build(story, onFirstPage=_add_page_number, onLaterPages=_add_page_number)


# ── WORD ACADÉMIQUE ───────────────────────────────────────────────────────────

def _set_doc_margins(doc, top=2.5, bottom=2.5, left=3.0, right=2.5):
    section = doc.sections[0]
    section.top_margin = Cm(top)
    section.bottom_margin = Cm(bottom)
    section.left_margin = Cm(left)
    section.right_margin = Cm(right)


def _add_word_header(doc, titre: str):
    header = doc.sections[0].header
    para = header.paragraphs[0]
    para.text = f"DevoirAI — {titre[:50]}"
    para.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    run = para.runs[0]
    run.font.size = Pt(9)
    run.font.color.rgb = DocxRGB(0x64, 0x74, 0x8b)


def _export_docx(text: str, path: str, titre: str):
    doc = Document()
    _set_doc_margins(doc)
    _add_word_header(doc, titre)

    # Page de titre
    title_para = doc.add_heading('', level=0)
    run = title_para.add_run(titre)
    run.font.size = Pt(22)
    run.font.color.rgb = DocxRGB(0x25, 0x63, 0xeb)
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

    date_para = doc.add_paragraph()
    date_run = date_para.add_run(datetime.now().strftime('%d/%m/%Y'))
    date_run.font.size = Pt(10)
    date_run.font.color.rgb = DocxRGB(0x64, 0x74, 0x8b)
    date_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_paragraph()

    # Ligne de séparation
    sep = doc.add_paragraph()
    sep_run = sep.add_run('─' * 60)
    sep_run.font.color.rgb = DocxRGB(0x25, 0x63, 0xeb)
    sep.alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_paragraph()

    # Contenu
    for line in text.split('\n'):
        line = line.strip()
        if not line:
            doc.add_paragraph()
            continue

        if line.startswith('# '):
            h = doc.add_heading(line[2:], level=1)
            h.runs[0].font.color.rgb = DocxRGB(0x1e, 0x40, 0xaf)
        elif line.startswith('## '):
            h = doc.add_heading(line[3:], level=2)
            h.runs[0].font.color.rgb = DocxRGB(0x7c, 0x3a, 0xed)
        elif line.startswith('### '):
            h = doc.add_heading(line[4:], level=3)
        elif line.startswith('**') and line.endswith('**'):
            p = doc.add_paragraph()
            run = p.add_run(line[2:-2])
            run.bold = True
            run.font.size = Pt(11)
        elif line.startswith('- ') or line.startswith('• '):
            p = doc.add_paragraph(line[2:], style='List Bullet')
            p.runs[0].font.size = Pt(11)
        else:
            p = doc.add_paragraph(line)
            p.runs[0].font.size = Pt(11)
            p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY

    doc.save(path)


# ── POWERPOINT THÈME VISUEL ────────────────────────────────────────────────────

PPTX_BG = RGBColor(0x0f, 0x17, 0x2a)
PPTX_PRIMARY = RGBColor(0x25, 0x63, 0xeb)
PPTX_ACCENT = RGBColor(0x7c, 0x3a, 0xed)
PPTX_WHITE = RGBColor(0xff, 0xff, 0xff)
PPTX_LIGHT = RGBColor(0xe2, 0xe8, 0xf0)
PPTX_MUTED = RGBColor(0x94, 0xa3, 0xb8)


def _set_slide_bg(slide, color: RGBColor):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text_box(slide, text, left, top, width, height,
                  font_size=18, bold=False, color=PPTX_WHITE,
                  align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(
        PInches(left), PInches(top), PInches(width), PInches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text[:500]
    run.font.size = PPt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def _add_accent_bar(slide, color=PPTX_PRIMARY):
    """Barre colorée en haut de la diapo."""
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        PInches(0), PInches(0),
        PInches(13.33), PInches(0.15)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


def _parse_slides(text: str) -> list[tuple[str, str]]:
    """Parse le format [SLIDE N] Titre\nContenu"""
    slide_pattern = re.compile(r'\[SLIDE\s*\d+\]\s*(.*)', re.IGNORECASE)
    slides_data = []
    current_title = ""
    current_content = []

    for line in text.split('\n'):
        match = slide_pattern.match(line.strip())
        if match:
            if current_title:
                slides_data.append((current_title, '\n'.join(current_content).strip()))
            current_title = match.group(1).strip()
            current_content = []
        elif current_title and line.strip():
            current_content.append(line.strip())

    if current_title:
        slides_data.append((current_title, '\n'.join(current_content).strip()))

    # Fallback si pas de format [SLIDE]
    if not slides_data:
        paras = [p.strip() for p in text.split('\n\n') if p.strip()]
        for para in paras[:12]:
            lines = para.split('\n')
            title = lines[0][:60].lstrip('#').strip()
            content = '\n'.join(lines[1:]) if len(lines) > 1 else para
            slides_data.append((title, content))

    return slides_data


def _export_pptx(text: str, path: str, titre: str):
    prs = Presentation()
    prs.slide_width = PInches(13.33)
    prs.slide_height = PInches(7.5)

    # ── Slide titre ──────────────────────────────────────────────────────────
    blank_layout = prs.slide_layouts[6]  # Blank
    title_slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(title_slide, PPTX_BG)

    # Gradient visuel - rectangle décoratif
    rect = title_slide.shapes.add_shape(
        1, PInches(0), PInches(2.5), PInches(13.33), PInches(3)
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0x1e, 0x40, 0xaf)
    rect.fill.fore_color.theme_color = None
    rect.line.fill.background()

    # Barre accent
    _add_accent_bar(title_slide, PPTX_PRIMARY)

    # Titre principal
    _add_text_box(
        title_slide, titre[:80],
        left=1, top=2.8, width=11.33, height=1.8,
        font_size=32, bold=True, color=PPTX_WHITE, align=PP_ALIGN.CENTER
    )

    # Sous-titre
    _add_text_box(
        title_slide,
        f"Présenté par — {datetime.now().strftime('%d/%m/%Y')}",
        left=1, top=4.8, width=11.33, height=0.5,
        font_size=14, color=PPTX_LIGHT, align=PP_ALIGN.CENTER, italic=True
    )

    # Logo DevoirAI
    _add_text_box(
        title_slide, "🎓 DevoirAI",
        left=0.5, top=6.7, width=3, height=0.5,
        font_size=11, color=PPTX_MUTED
    )

    # ── Slides contenu ────────────────────────────────────────────────────────
    slides_data = _parse_slides(text)
    colors_rotation = [PPTX_PRIMARY, PPTX_ACCENT, RGBColor(0x05, 0x96, 0x69)]

    for i, (slide_title, slide_content) in enumerate(slides_data):
        slide = prs.slides.add_slide(blank_layout)
        _set_slide_bg(slide, PPTX_BG)

        accent_color = colors_rotation[i % len(colors_rotation)]
        _add_accent_bar(slide, accent_color)

        # Numéro de slide
        _add_text_box(
            slide, f"{i + 1:02d}",
            left=11.8, top=0.2, width=1.2, height=0.6,
            font_size=11, color=PPTX_MUTED, align=PP_ALIGN.RIGHT
        )

        # Titre slide
        _add_text_box(
            slide, slide_title[:70],
            left=0.5, top=0.3, width=11, height=0.9,
            font_size=24, bold=True, color=PPTX_WHITE
        )

        # Ligne décorative sous le titre
        line_shape = slide.shapes.add_shape(
            1, PInches(0.5), PInches(1.35), PInches(2), PInches(0.04)
        )
        line_shape.fill.solid()
        line_shape.fill.fore_color.rgb = accent_color
        line_shape.line.fill.background()

        # Contenu
        _add_text_box(
            slide, slide_content[:700],
            left=0.5, top=1.6, width=12.3, height=5.4,
            font_size=16, color=PPTX_LIGHT
        )

    prs.save(path)
